from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── Brand Palette ─────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1A, 0x29, 0x4E)   # deep navy
BLUE      = RGBColor(0x1F, 0x49, 0x7D)   # brand blue
ORANGE    = RGBColor(0xE8, 0x6A, 0x2C)   # brand orange
ORANGE2   = RGBColor(0xFF, 0x8C, 0x00)   # accent orange
TEAL      = RGBColor(0x2E, 0x86, 0xAB)   # teal accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY     = RGBColor(0xF0, 0xF4, 0xF8)
MGRAY     = RGBColor(0x9E, 0x9E, 0x9E)
DGRAY     = RGBColor(0x37, 0x37, 0x37)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
RED       = RGBColor(0xC0, 0x39, 0x2B)
GOLD      = RGBColor(0xF3, 0x9C, 0x12)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # truly blank

# ── Low-level helpers ─────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

def rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=18, bold=False, italic=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def mtxt(slide, lines, x, y, w, h,
         size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
         spacing=1.15, font="Calibri"):
    """Multi-line textbox where each item in lines is (text, bold_override, size_override, color_override)."""
    from pptx.util import Pt as pt2
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            t, b, s, c = item, bold, size, color
        else:
            t = item[0]
            b = item[1] if len(item) > 1 else bold
            s = item[2] if len(item) > 2 else size
            c = item[3] if len(item) > 3 else color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(spacing * 100000)))
        run = p.add_run()
        run.text = t
        run.font.name  = font
        run.font.size  = Pt(s)
        run.font.bold  = b
        run.font.color.rgb = c
    return tb

def circle(slide, cx, cy, r, fill, line=None):
    from pptx.util import Emu
    sh = slide.shapes.add_shape(9, cx - r, cy - r, r*2, r*2)  # OVAL
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
    else:
        sh.line.fill.background()
    return sh

def add_slide():
    return prs.slides.add_slide(BLANK)

def gradient_bg(slide, top_color, bot_color):
    """Simulate gradient with two overlapping rects."""
    rect(slide, 0, 0, W, H, fill=top_color)
    # overlay a semi-transparent feel with a slightly lighter band
    rect(slide, 0, H*0.55, W, H*0.45, fill=bot_color)

def pill(slide, x, y, w, h, fill, text, tsize=14, tcol=WHITE, tbold=True):
    from pptx.util import Emu
    sh = slide.shapes.add_shape(5, x, y, w, h)   # ROUNDED_RECTANGLE
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.adjustments[0] = 0.5
    tf = sh.text_frame; tf.word_wrap = False
    p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.name  = "Calibri"; run.font.size = Pt(tsize)
    run.font.bold  = tbold;     run.font.color.rgb = tcol
    return sh

def icon_card(slide, x, y, w, h, icon, title, body_lines, bg=BLUE, icon_col=ORANGE):
    rect(slide, x, y, w, h, fill=bg)
    # top colour band
    rect(slide, x, y, w, Inches(0.06), fill=icon_col)
    txt(slide, icon, x, y + Inches(0.18), w, Inches(0.55),
        size=28, bold=True, color=icon_col, align=PP_ALIGN.CENTER)
    txt(slide, title, x, y + Inches(0.72), w, Inches(0.38),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    body_text = "\n".join(body_lines)
    txt(slide, body_text, x + Inches(0.15), y + Inches(1.08),
        w - Inches(0.3), h - Inches(1.15),
        size=10.5, color=RGBColor(0xCC, 0xD6, 0xE8), align=PP_ALIGN.LEFT)

def section_label(slide, text):
    pill(slide, Inches(0.35), Inches(0.22), Inches(2.6), Inches(0.32),
         fill=ORANGE, text=text, tsize=11)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=NAVY)
# diagonal accent band
from pptx.util import Emu
import math

# right orange accent strip
rect(sl, W - Inches(3.8), 0, Inches(3.8), H, fill=rgb(0x1F, 0x3A, 0x6A))
rect(sl, W - Inches(0.55), 0, Inches(0.55), H, fill=ORANGE)

# large logo text
txt(sl, "Shiftzy", Inches(0.6), Inches(1.5), Inches(6), Inches(1.4),
    size=72, bold=True, color=WHITE, font="Calibri")
txt(sl, "Go", Inches(0.6), Inches(2.75), Inches(4), Inches(1.2),
    size=72, bold=True, color=ORANGE, font="Calibri")

# tagline
txt(sl, "India's Smart Vehicle Shifting &\nTravel Cost-Sharing Platform",
    Inches(0.6), Inches(4.0), Inches(7.5), Inches(1.0),
    size=22, italic=True, color=rgb(0xB0, 0xC4, 0xDE))

# divider line
rect(sl, Inches(0.6), Inches(5.1), Inches(5.5), Inches(0.04), fill=ORANGE)

# sub-line
txt(sl, "Investor & Stakeholder Presentation  ·  July 2026",
    Inches(0.6), Inches(5.25), Inches(8), Inches(0.45),
    size=14, color=rgb(0x90, 0xA8, 0xC8))

# right side stats
for i, (num, label) in enumerate([("4,370+", "Trusted Users"),
                                    ("50+", "Routes Live"),
                                    ("₹8K→₹4K", "Cost Saving")]):
    bx = Inches(9.8)
    by = Inches(1.5 + i * 1.8)
    txt(sl, num,  bx, by,          Inches(3.0), Inches(0.75),
        size=36, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txt(sl, label, bx, by + Inches(0.65), Inches(3.0), Inches(0.4),
        size=13, color=rgb(0xB0, 0xC4, 0xDE), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=LGRAY)
rect(sl, 0, 0, Inches(4.4), H, fill=NAVY)
rect(sl, 0, 0, Inches(0.08), H, fill=ORANGE)

txt(sl, "What We'll\nCover Today",
    Inches(0.25), Inches(0.5), Inches(4.0), Inches(1.8),
    size=30, bold=True, color=WHITE)

items = [
    ("01", "The Problem"),
    ("02", "Our Solution"),
    ("03", "How It Works"),
    ("04", "Market Opportunity"),
    ("05", "Business Model"),
    ("06", "Insurance Partnership"),
    ("07", "Financial Projections"),
    ("08", "Growth Roadmap"),
    ("09", "Why Invest"),
    ("10", "Next Steps"),
]
for i, (num, label) in enumerate(items):
    col = 0 if i < 5 else 1
    row = i % 5
    bx = Inches(4.8 + col * 4.2)
    by = Inches(0.6 + row * 1.2)
    rect(sl, bx, by, Inches(3.8), Inches(0.9), fill=WHITE)
    rect(sl, bx, by, Inches(0.06), Inches(0.9), fill=ORANGE)
    txt(sl, num, bx + Inches(0.18), by + Inches(0.1), Inches(0.6), Inches(0.7),
        size=22, bold=True, color=ORANGE)
    txt(sl, label, bx + Inches(0.75), by + Inches(0.2), Inches(3.0), Inches(0.55),
        size=15, bold=True, color=DGRAY)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=LGRAY)
rect(sl, 0, 0, W, Inches(1.3), fill=NAVY)
rect(sl, 0, 0, W, Inches(0.07), fill=ORANGE)

section_label(sl, "01  THE PROBLEM")
txt(sl, "Every Day, Indians Face Two Expensive Challenges",
    Inches(0.4), Inches(0.28), Inches(12), Inches(0.85),
    size=28, bold=True, color=WHITE)

# Card 1
rect(sl, Inches(0.35), Inches(1.5), Inches(5.9), Inches(5.2), fill=WHITE)
rect(sl, Inches(0.35), Inches(1.5), Inches(5.9), Inches(0.07), fill=RED)
txt(sl, "🚗", Inches(0.35), Inches(1.55), Inches(5.9), Inches(0.9),
    size=40, align=PP_ALIGN.CENTER, color=DGRAY)
txt(sl, "VEHICLE OWNERS",
    Inches(0.5), Inches(2.42), Inches(5.6), Inches(0.5),
    size=17, bold=True, color=RED, align=PP_ALIGN.CENTER)
txt(sl, "Need to move their car, bike, or SUV\nto another city — but options are painful:",
    Inches(0.5), Inches(2.9), Inches(5.6), Inches(0.7),
    size=12.5, color=DGRAY, align=PP_ALIGN.CENTER)
for item in [
    "🚛  Truck/trailer carrier → ₹8,000 – ₹30,000",
    "👤  Hire a stranger driver → Safety risk",
    "🛣️  Drive it yourself → Costs time & money",
]:
    txt(sl, item, Inches(0.6), Inches(3.6 + [item].index(item) * 0.0), Inches(5.4), Inches(0.45),
        size=12, color=DGRAY)

# manually place the bullet items
for idx, item in enumerate([
    "🚛  Truck/trailer carrier → ₹8,000 – ₹30,000",
    "👤  Hire a stranger driver → Safety risk",
    "🛣️  Drive it yourself → Costs time & money",
    "😓  No trusted, affordable alternative exists",
]):
    txt(sl, item, Inches(0.65), Inches(3.58 + idx * 0.58), Inches(5.4), Inches(0.52),
        size=12, color=DGRAY)

# Card 2
rect(sl, Inches(6.65), Inches(1.5), Inches(5.9), Inches(5.2), fill=WHITE)
rect(sl, Inches(6.65), Inches(1.5), Inches(5.9), Inches(0.07), fill=BLUE)
txt(sl, "✈️", Inches(6.65), Inches(1.55), Inches(5.9), Inches(0.9),
    size=40, align=PP_ALIGN.CENTER, color=DGRAY)
txt(sl, "TRAVELERS",
    Inches(6.8), Inches(2.42), Inches(5.6), Inches(0.5),
    size=17, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
txt(sl, "Pay the full cost of intercity travel —\nfuel, tolls, food, tickets — alone:",
    Inches(6.8), Inches(2.9), Inches(5.6), Inches(0.7),
    size=12.5, color=DGRAY, align=PP_ALIGN.CENTER)
for idx, item in enumerate([
    "⛽  Full fuel + toll → ₹2,500 – ₹10,000",
    "🎫  Train/bus + last-mile taxi → ₹1,800+",
    "⏱️  Long journey, high cost, zero benefit",
    "💸  No way to offset your travel cost",
]):
    txt(sl, item, Inches(6.95), Inches(3.58 + idx * 0.58), Inches(5.4), Inches(0.52),
        size=12, color=DGRAY)

# bottom stat
rect(sl, Inches(2.5), Inches(6.55), Inches(8.3), Inches(0.72), fill=NAVY)
txt(sl, "₹50,000 Crore+  lost every year across India in unoptimised vehicle transport & solo travel costs",
    Inches(2.6), Inches(6.6), Inches(8.1), Inches(0.6),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=NAVY)
rect(sl, 0, 0, W, Inches(0.07), fill=ORANGE)

section_label(sl, "02  OUR SOLUTION")
txt(sl, "Introducing Shiftzy Go",
    Inches(0.4), Inches(0.28), Inches(12), Inches(0.85),
    size=32, bold=True, color=WHITE)

# Big tagline
txt(sl, "One App. Two Winners. Half the Cost.",
    Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.7),
    size=24, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# Solution visual
rect(sl, Inches(0.4), Inches(1.9), Inches(5.5), Inches(4.6), fill=rgb(0x1F, 0x3A, 0x6A))
txt(sl, "🚗\nVEHICLE OWNER",
    Inches(0.5), Inches(2.0), Inches(5.3), Inches(1.0),
    size=15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txt(sl, "Posts vehicle shift request\n\nPickup & drop city\n\nVehicle details + photos\n\nPreferred dates",
    Inches(0.6), Inches(3.1), Inches(5.1), Inches(2.0),
    size=12.5, color=rgb(0xCC, 0xD6, 0xE8))
txt(sl, "✅  Pays only 50% of fuel + toll",
    Inches(0.6), Inches(5.1), Inches(5.1), Inches(0.45),
    size=12, bold=True, color=GREEN)
txt(sl, "✅  Vehicle safely shifted by verified driver",
    Inches(0.6), Inches(5.55), Inches(5.1), Inches(0.45),
    size=12, bold=True, color=GREEN)

# arrow in middle
txt(sl, "🤝\nSHIFTZY GO\nMATCHES THEM",
    Inches(5.9), Inches(2.8), Inches(1.55), Inches(1.8),
    size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

rect(sl, Inches(7.45), Inches(1.9), Inches(5.5), Inches(4.6), fill=rgb(0x1F, 0x3A, 0x6A))
txt(sl, "🧳\nTRAVELER",
    Inches(7.55), Inches(2.0), Inches(5.3), Inches(1.0),
    size=15, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
txt(sl, "Browses open shift requests\n\nApplies to drive on same route\n\nVerified via KYC + DL check\n\nPayment through app",
    Inches(7.65), Inches(3.1), Inches(5.1), Inches(2.0),
    size=12.5, color=rgb(0xCC, 0xD6, 0xE8))
txt(sl, "✅  Pays only 50% of fuel + toll + small fee",
    Inches(7.65), Inches(5.1), Inches(5.1), Inches(0.45),
    size=12, bold=True, color=GREEN)
txt(sl, "✅  Reaches destination at near-zero cost",
    Inches(7.65), Inches(5.55), Inches(5.1), Inches(0.45),
    size=12, bold=True, color=GREEN)

# bottom bar
rect(sl, 0, Inches(6.72), W, Inches(0.78), fill=ORANGE)
txt(sl, "RESULT: Vehicle shifted ✦ Traveler travels ✦ Both pay HALF ✦ Everyone wins",
    Inches(0.3), Inches(6.78), Inches(12.7), Inches(0.6),
    size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — HOW IT WORKS
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=LGRAY)
rect(sl, 0, 0, W, Inches(1.3), fill=NAVY)
rect(sl, 0, 0, W, Inches(0.07), fill=ORANGE)

section_label(sl, "03  HOW IT WORKS")
txt(sl, "Simple. Safe. Smart. — 6 Steps to a Successful Shift",
    Inches(0.4), Inches(0.28), Inches(12.5), Inches(0.85),
    size=26, bold=True, color=WHITE)

steps = [
    ("01", "OWNER POSTS", "🚗", "Owner lists their vehicle, cities, dates & photos on the app", BLUE),
    ("02", "TRAVELER APPLIES", "🙋", "Verified traveler heading same route applies to drive", TEAL),
    ("03", "OWNER APPROVES", "✅", "Owner reviews profile, rating & DL — approves the match", ORANGE),
    ("04", "PAYMENT & INSURANCE", "🔒", "Traveler pays their 50% share. Insurance auto-activates", rgb(0x8E, 0x44, 0xAD)),
    ("05", "VEHICLE SHIFTS", "🛣️", "Traveler drives, GPS tracked, both get live updates", GREEN),
    ("06", "DELIVERY & REVIEW", "⭐", "Owner confirms delivery. Both rate each other. Done!", GOLD),
]

for i, (num, title, icon, desc, col) in enumerate(steps):
    col_i = i % 3
    row_i = i // 3
    bx = Inches(0.3 + col_i * 4.35)
    by = Inches(1.45 + row_i * 2.8)
    rect(sl, bx, by, Inches(4.1), Inches(2.55), fill=WHITE)
    rect(sl, bx, by, Inches(4.1), Inches(0.06), fill=col)
    # number badge
    sh = circle(sl, bx + Inches(0.52), by + Inches(0.52), Inches(0.36), col)
    txt(sl, num, bx + Inches(0.17), by + Inches(0.15), Inches(0.7), Inches(0.72),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, icon, bx + Inches(0.85), by + Inches(0.1), Inches(2.8), Inches(0.72),
        size=30, color=DGRAY)
    txt(sl, title, bx + Inches(0.2), by + Inches(0.78), Inches(3.7), Inches(0.45),
        size=13, bold=True, color=col)
    txt(sl, desc, bx + Inches(0.2), by + Inches(1.2), Inches(3.7), Inches(1.2),
        size=11, color=DGRAY)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — THE COST-SHARING MAGIC
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=NAVY)
rect(sl, 0, 0, W, Inches(0.07), fill=ORANGE)

section_label(sl, "04  THE COST-SHARING MAGIC")
txt(sl, "The Numbers That Make Everyone Smile",
    Inches(0.4), Inches(0.28), Inches(12), Inches(0.85),
    size=28, bold=True, color=WHITE)

# Example route
rect(sl, Inches(0.35), Inches(1.18), Inches(12.63), Inches(0.68), fill=rgb(0x1F, 0x3A, 0x6A))
txt(sl, "📍  Chennai  →  Bangalore  ·  Sedan  ·  346 km",
    Inches(0.5), Inches(1.22), Inches(12.3), Inches(0.58),
    size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# Without Shiftzy
rect(sl, Inches(0.35), Inches(2.05), Inches(5.7), Inches(4.75), fill=rgb(0x2D, 0x1A, 0x1A))
rect(sl, Inches(0.35), Inches(2.05), Inches(5.7), Inches(0.06), fill=RED)
txt(sl, "❌  WITHOUT SHIFTZY GO",
    Inches(0.45), Inches(2.12), Inches(5.5), Inches(0.5),
    size=14, bold=True, color=RED, align=PP_ALIGN.CENTER)

without_rows = [
    ("Fuel (Sedan, 18 km/l)", "₹2,697"),
    ("Toll charges", "₹519"),
    ("Total cost if alone", "₹3,216"),
    ("Owner pays alone", "₹3,216"),
    ("Traveler cost (bus+taxi)", "₹2,500+"),
    ("COMBINED SPEND", "₹5,716+"),
]
for idx, (label, val) in enumerate(without_rows):
    by = Inches(2.72 + idx * 0.57)
    is_total = "COMBINED" in label or "Total" in label
    rect(sl, Inches(0.45), by, Inches(5.5), Inches(0.5),
         fill=rgb(0x40, 0x20, 0x20) if is_total else rgb(0x38, 0x22, 0x22))
    txt(sl, label, Inches(0.6), by + Inches(0.07), Inches(3.8), Inches(0.4),
        size=11.5 if not is_total else 13,
        bold=is_total, color=rgb(0xFF, 0xCC, 0xCC) if not is_total else RED)
    txt(sl, val,  Inches(4.4), by + Inches(0.07), Inches(1.4), Inches(0.4),
        size=11.5 if not is_total else 13,
        bold=is_total, color=rgb(0xFF, 0xCC, 0xCC) if not is_total else RED,
        align=PP_ALIGN.RIGHT)

# VS label
txt(sl, "VS", Inches(6.1), Inches(4.0), Inches(1.1), Inches(0.9),
    size=36, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# With Shiftzy
rect(sl, Inches(7.3), Inches(2.05), Inches(5.7), Inches(4.75), fill=rgb(0x12, 0x2A, 0x1A))
rect(sl, Inches(7.3), Inches(2.05), Inches(5.7), Inches(0.06), fill=GREEN)
txt(sl, "✅  WITH SHIFTZY GO",
    Inches(7.4), Inches(2.12), Inches(5.5), Inches(0.5),
    size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

with_rows = [
    ("Full fuel + toll (shared)", "₹3,216"),
    ("Owner's 50% share", "₹1,608"),
    ("Traveler's 50% share", "₹1,608"),
    ("Platform fee + GST", "₹190"),
    ("Traveler total payable", "₹1,798"),
    ("COMBINED SPEND", "₹3,406"),
]
for idx, (label, val) in enumerate(with_rows):
    by = Inches(2.72 + idx * 0.57)
    is_total = "COMBINED" in label
    rect(sl, Inches(7.4), by, Inches(5.5), Inches(0.5),
         fill=rgb(0x1A, 0x40, 0x28) if is_total else rgb(0x16, 0x32, 0x20))
    txt(sl, label, Inches(7.55), by + Inches(0.07), Inches(3.8), Inches(0.4),
        size=11.5 if not is_total else 13,
        bold=is_total, color=rgb(0xCC, 0xFF, 0xCC) if not is_total else GREEN)
    txt(sl, val,  Inches(11.35), by + Inches(0.07), Inches(1.4), Inches(0.4),
        size=11.5 if not is_total else 13,
        bold=is_total, color=rgb(0xCC, 0xFF, 0xCC) if not is_total else GREEN,
        align=PP_ALIGN.RIGHT)

# savings banner
rect(sl, Inches(0.35), Inches(6.72), Inches(12.63), Inches(0.65), fill=ORANGE)
txt(sl, "🎉  Combined saving: ₹2,310 per trip  ·  Owner saves ₹1,608  ·  Traveler saves ₹702 vs bus+taxi",
    Inches(0.4), Inches(6.77), Inches(12.5), Inches(0.55),
    size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — MARKET OPPORTUNITY
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=LGRAY)
rect(sl, 0, 0, W, Inches(1.3), fill=NAVY)
rect(sl, 0, 0, W, Inches(0.07), fill=ORANGE)

section_label(sl, "05  MARKET OPPORTUNITY")
txt(sl, "A Massive, Untapped Market Waiting to Be Disrupted",
    Inches(0.4), Inches(0.28), Inches(12.5), Inches(0.85),
    size=26, bold=True, color=WHITE)

stats = [
    ("30 Crore+", "Registered vehicles\nin India", BLUE),
    ("₹85,000 Cr", "Annual vehicle logistics\nmarket size", ORANGE),
    ("3 Crore+", "Intercity trips\nper month", TEAL),
    ("82%", "Users open to cost-sharing\nplatforms (survey)", GREEN),
]
for i, (num, label, col) in enumerate(stats):
    bx = Inches(0.3 + i * 3.25)
    rect(sl, bx, Inches(1.45), Inches(3.0), Inches(2.1), fill=WHITE)
    rect(sl, bx, Inches(1.45), Inches(3.0), Inches(0.07), fill=col)
    txt(sl, num, bx, Inches(1.55), Inches(3.0), Inches(0.85),
        size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, label, bx, Inches(2.38), Inches(3.0), Inches(0.75),
        size=12, color=DGRAY, align=PP_ALIGN.CENTER)

# TAM SAM SOM
txt(sl, "Market Sizing — TAM · SAM · SOM",
    Inches(0.4), Inches(3.78), Inches(7), Inches(0.5),
    size=16, bold=True, color=NAVY)
market_rows = [
    ("TAM — Total Addressable Market", "All intercity vehicle transport & travel in India", "₹85,000 Cr", BLUE),
    ("SAM — Serviceable Addressable Market", "Peer-to-peer vehicle shifting corridor (Tier 1 + 2 cities)", "₹12,000 Cr", TEAL),
    ("SOM — Serviceable Obtainable Market", "Shiftzy Go's realistic 3-year capture (0.3% of SAM)", "₹360 Cr", ORANGE),
]
for idx, (seg, desc, val, col) in enumerate(market_rows):
    by = Inches(4.35 + idx * 0.9)
    rect(sl, Inches(0.35), by, Inches(9.0), Inches(0.78), fill=WHITE)
    rect(sl, Inches(0.35), by, Inches(0.07), Inches(0.78), fill=col)
    txt(sl, seg,  Inches(0.55), by + Inches(0.08), Inches(4.5), Inches(0.35),
        size=12, bold=True, color=col)
    txt(sl, desc, Inches(0.55), by + Inches(0.42), Inches(5.5), Inches(0.3),
        size=10.5, color=DGRAY)
    txt(sl, val,  Inches(8.0), by + Inches(0.18), Inches(1.2), Inches(0.45),
        size=16, bold=True, color=col, align=PP_ALIGN.RIGHT)

# right side: why now
rect(sl, Inches(9.8), Inches(3.68), Inches(3.2), Inches(3.6), fill=NAVY)
txt(sl, "WHY NOW?", Inches(9.9), Inches(3.78), Inches(3.0), Inches(0.45),
    size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
for idx, why in enumerate([
    "📱  India's UPI/digital\npayment boom",
    "🚦  Vehicle ownership\nup 12% YoY",
    "💼  Post-pandemic\nremote work = relocation",
    "🌿  Sustainability push\n= shared transport",
    "📍  No platform exists\nfor this exact gap",
]):
    txt(sl, why, Inches(9.95), Inches(4.32 + idx * 0.57), Inches(2.9), Inches(0.5),
        size=10.5, color=rgb(0xCC, 0xD6, 0xE8))

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — BUSINESS MODEL
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=NAVY)
rect(sl, 0, 0, W, Inches(0.07), fill=ORANGE)

section_label(sl, "06  BUSINESS MODEL")
txt(sl, "Multiple Revenue Streams — Built to Scale",
    Inches(0.4), Inches(0.28), Inches(12), Inches(0.85),
    size=28, bold=True, color=WHITE)

streams = [
    ("💰", "Platform Fee", "10% of traveler's\n50% share per trip.\nCore revenue stream.", ORANGE),
    ("🛡️", "Insurance Commission", "12–15% distribution\nfee on every per-trip\ninsurance premium.", TEAL),
    ("⭐", "Premium Listings", "Vehicle owners pay\nfor featured placement\n& priority matching.", rgb(0x8E, 0x44, 0xAD)),
    ("🏢", "Corporate Accounts", "Company vehicle\nshifts at bulk rates\n(B2B segment).", BLUE),
    ("📊", "Data & Analytics", "Fleet insights sold\nto OEMs, insurers &\nlogistics companies.", GREEN),
    ("🎯", "Ads & Partnerships", "Fuel brands, service\ncentres, toll operators\n& travel partners.", GOLD),
]
for i, (icon, title, desc, col) in enumerate(streams):
    col_i = i % 3
    row_i = i // 3
    bx = Inches(0.35 + col_i * 4.35)
    by = Inches(1.35 + row_i * 2.7)
    rect(sl, bx, by, Inches(4.1), Inches(2.45), fill=rgb(0x1F, 0x3A, 0x6A))
    rect(sl, bx, by, Inches(4.1), Inches(0.06), fill=col)
    txt(sl, icon, bx, by + Inches(0.1), Inches(4.1), Inches(0.65),
        size=30, align=PP_ALIGN.CENTER, color=WHITE)
    txt(sl, title, bx, by + Inches(0.72), Inches(4.1), Inches(0.45),
        size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, desc, bx + Inches(0.25), by + Inches(1.15), Inches(3.6), Inches(1.2),
        size=11.5, color=rgb(0xB0, 0xC4, 0xDE), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — INSURANCE PARTNERSHIP
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=LGRAY)
rect(sl, 0, 0, W, Inches(1.3), fill=NAVY)
rect(sl, 0, 0, W, Inches(0.07), fill=ORANGE)

section_label(sl, "07  INSURANCE PARTNERSHIP")
txt(sl, "Safety First — Insurance is the Backbone of Shiftzy Go",
    Inches(0.4), Inches(0.28), Inches(12.5), Inches(0.85),
    size=24, bold=True, color=WHITE)

# 4 coverage cards
covers = [
    ("🚘", "Vehicle Transit\nInsurance", "Covers physical damage to owner's vehicle from trip start to delivery", BLUE),
    ("⚖️", "Third-Party\nLiability", "Covers bodily injury & property damage to third parties (MV Act 1988)", ORANGE),
    ("👷", "Personal Accident\nCover — Traveler", "₹10 lakh accidental death/disability cover for the traveler during trip", TEAL),
    ("📋", "Owner Legal\nLiability Cover", "Protects vehicle owner from legal liability from traveler's actions", rgb(0x8E, 0x44, 0xAD)),
]
for i, (icon, title, desc, col) in enumerate(covers):
    bx = Inches(0.3 + i * 3.25)
    rect(sl, bx, Inches(1.5), Inches(3.0), Inches(3.3), fill=WHITE)
    rect(sl, bx, Inches(1.5), Inches(3.0), Inches(0.07), fill=col)
    txt(sl, icon, bx, Inches(1.58), Inches(3.0), Inches(0.75),
        size=32, align=PP_ALIGN.CENTER, color=DGRAY)
    txt(sl, title, bx + Inches(0.1), Inches(2.3), Inches(2.8), Inches(0.7),
        size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, desc, bx + Inches(0.15), Inches(3.0), Inches(2.7), Inches(1.65),
        size=11, color=DGRAY, align=PP_ALIGN.CENTER)

# Partnership benefits
txt(sl, "Why an Insurer Should Partner With Us",
    Inches(0.4), Inches(4.98), Inches(8), Inches(0.5),
    size=16, bold=True, color=NAVY)

benefits = [
    ("🥇 First-Mover", "No existing product covers peer-to-peer vehicle transit in India — own this category exclusively"),
    ("📈 Scale", "₹15–25 lakh premium Year 1 → ₹3–5 crore by Year 3 → growing automatically with platform"),
    ("🔒 Low Claims Risk", "Pre-trip photos, GPS tracking, DL verification & rating system significantly reduce fraud"),
    ("💎 Premium Demographic", "Users aged 21–40, digital-native, creditworthy — most valuable insurance segment"),
]
for i, (title, desc) in enumerate(benefits):
    col_i = i % 2
    row_i = i // 2
    bx = Inches(0.3 + col_i * 6.5)
    by = Inches(5.55 + row_i * 0.72)
    rect(sl, bx, by, Inches(6.2), Inches(0.62), fill=WHITE)
    rect(sl, bx, by, Inches(0.06), Inches(0.62), fill=ORANGE)
    txt(sl, title, bx + Inches(0.18), by + Inches(0.1), Inches(1.6), Inches(0.42),
        size=11, bold=True, color=BLUE)
    txt(sl, desc, bx + Inches(1.78), by + Inches(0.1), Inches(4.3), Inches(0.42),
        size=10.5, color=DGRAY)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — PLATFORM FEATURES
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=NAVY)
rect(sl, 0, 0, W, Inches(0.07), fill=ORANGE)

section_label(sl, "08  PLATFORM FEATURES")
txt(sl, "A Complete, Technology-First Platform",
    Inches(0.4), Inches(0.28), Inches(12), Inches(0.85),
    size=28, bold=True, color=WHITE)

features = [
    ("📱", "Mobile-First App",    "Intuitive iOS & Android experience with seamless booking flow"),
    ("🔍", "Smart Matching",      "AI-powered route & date matching between owners and travelers"),
    ("✅", "KYC Verification",    "Aadhaar + DL verification before any traveler can drive"),
    ("📸", "Pre-Trip Photos",     "Time-stamped, geo-tagged inspection photos for every trip"),
    ("🗺️", "Live GPS Tracking",  "Real-time route monitoring for both parties during transit"),
    ("💬", "In-App Chat",         "Direct, encrypted messaging between owner and traveler"),
    ("💳", "Secure Payments",     "Stripe-powered payments with transparent fee breakdown"),
    ("⭐", "Dual Rating System",  "5-star ratings for both travelers and vehicle owners"),
    ("🛡️", "Auto-Insurance",     "Insurance policy auto-generated and delivered at booking"),
    ("🤖", "AI Chat Support",     "24/7 intelligent FAQ bot with human escalation path"),
    ("📊", "Admin Dashboard",     "Real-time operational metrics, fraud monitoring & controls"),
    ("🔔", "Smart Notifications", "Push alerts for trip status, approvals & messages"),
]
for i, (icon, title, desc) in enumerate(features):
    col_i = i % 4
    row_i = i // 4
    bx = Inches(0.3 + col_i * 3.26)
    by = Inches(1.35 + row_i * 1.9)
    rect(sl, bx, by, Inches(3.05), Inches(1.72), fill=rgb(0x1F, 0x3A, 0x6A))
    rect(sl, bx, by, Inches(3.05), Inches(0.05), fill=ORANGE if i % 3 == 0 else (TEAL if i % 3 == 1 else BLUE))
    txt(sl, icon + "  " + title,
        bx + Inches(0.15), by + Inches(0.1), Inches(2.75), Inches(0.5),
        size=12, bold=True, color=WHITE)
    txt(sl, desc,
        bx + Inches(0.15), by + Inches(0.58), Inches(2.75), Inches(1.05),
        size=10, color=rgb(0xB0, 0xC4, 0xDE))

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — FINANCIAL PROJECTIONS
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=LGRAY)
rect(sl, 0, 0, W, Inches(1.3), fill=NAVY)
rect(sl, 0, 0, W, Inches(0.07), fill=ORANGE)

section_label(sl, "09  FINANCIAL PROJECTIONS")
txt(sl, "Conservative, Milestone-Based Growth",
    Inches(0.4), Inches(0.28), Inches(12), Inches(0.85),
    size=26, bold=True, color=WHITE)

# Main projections table
headers = ["Metric", "Year 1", "Year 2", "Year 3"]
rows = [
    ["Active Cities",        "7 (South India)",   "20 (Pan India)",    "35+ cities"],
    ["Monthly Trips",        "500 – 800",          "2,500 – 4,000",     "10,000+"],
    ["Annual Trips",         "~8,000",             "~40,000",           "~1,50,000"],
    ["Registered Users",     "25,000",             "1,50,000",          "5,00,000+"],
    ["Platform Fee Revenue", "₹60 – 80 Lakh",      "₹3 – 4.5 Crore",   "₹12 – 20 Crore"],
    ["Insurance Premium",    "₹15 – 25 Lakh",      "₹75L – 1.2 Crore", "₹3 – 5 Crore"],
    ["Total Revenue",        "₹75 Lakh – 1 Crore", "₹4 – 6 Crore",     "₹15 – 25 Crore"],
    ["Net Margin (est.)",    "–15% (invest phase)","12 – 18%",          "22 – 28%"],
]
col_colors = [NAVY, BLUE, TEAL, ORANGE]
col_w = [2.6, 2.6, 2.6, 2.6]
bx0 = Inches(0.35)
by0 = Inches(1.45)
# header row
for ci, (h, col) in enumerate(zip(headers, col_colors)):
    bx = bx0 + Inches(ci * 2.6)
    rect(sl, bx, by0, Inches(2.58), Inches(0.52), fill=col)
    txt(sl, h, bx, by0 + Inches(0.09), Inches(2.58), Inches(0.38),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for ri, row in enumerate(rows):
    by = by0 + Inches(0.55 + ri * 0.62)
    fill_bg = WHITE if ri % 2 == 0 else rgb(0xF0, 0xF4, 0xF8)
    for ci, val in enumerate(row):
        bx = bx0 + Inches(ci * 2.6)
        rect(sl, bx, by, Inches(2.58), Inches(0.59), fill=fill_bg)
        is_rev = "Revenue" in row[0] or "Margin" in row[0]
        tcol = BLUE if ci > 0 and is_rev else DGRAY
        tbold = is_rev and ci > 0
        txt(sl, val, bx + Inches(0.1), by + Inches(0.1), Inches(2.38), Inches(0.42),
            size=11 if ci > 0 else 11.5, bold=tbold or ci == 0,
            color=tcol if ci > 0 else NAVY, align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

# key assumptions
rect(sl, Inches(10.9), Inches(1.45), Inches(2.1), Inches(5.82), fill=NAVY)
txt(sl, "KEY\nASSUMPTIONS", Inches(10.95), Inches(1.55), Inches(2.0), Inches(0.7),
    size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
for idx, note in enumerate([
    "Avg. trip value\n₹1,800 – ₹3,500",
    "Platform fee\n10% of traveler\nshare",
    "Insurance comm.\n12% of premium",
    "CAC (cost/user)\n₹120 Y1 → ₹60 Y3",
    "Churn <8%\nmonth-on-month",
    "B2B corporate\nfrom Month 14",
]):
    txt(sl, note, Inches(10.95), Inches(2.32 + idx * 0.72), Inches(2.0), Inches(0.68),
        size=9.5, color=rgb(0xCC, 0xD6, 0xE8), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — GROWTH ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=NAVY)
rect(sl, 0, 0, W, Inches(0.07), fill=ORANGE)

section_label(sl, "10  GROWTH ROADMAP")
txt(sl, "Four Phases. One Vision — India's #1 Vehicle Shifting Platform",
    Inches(0.4), Inches(0.28), Inches(12.5), Inches(0.85),
    size=24, bold=True, color=WHITE)

phases = [
    ("PHASE 1", "Months\n1 – 12", "🌱 South India Launch",
     ["Chennai, Bangalore,\nHyderabad, Coimbatore,\nMadurai, Tirunelveli, Kochi",
      "500–800 trips/month",
      "Basic transit insurance",
      "25,000 users"],
     ORANGE),
    ("PHASE 2", "Months\n12 – 24", "🚀 Pan-India Expansion",
     ["Mumbai, Pune, Delhi NCR,\nAhmedabad, Jaipur, Kolkata",
      "2,500–4,000 trips/month",
      "Corporate fleet accounts",
      "1,50,000 users"],
     TEAL),
    ("PHASE 3", "Months\n24 – 36", "⚡ Platform Ecosystem",
     ["All metro + Tier 2 cities",
      "OEM dealer partnerships",
      "Annual subscription insurance",
      "5,00,000 users"],
     BLUE),
    ("PHASE 4", "Years\n3 – 5", "🏆 Market Leadership",
     ["30+ cities + int'l corridors",
      "10,000+ trips/month",
      "IPO readiness",
      "10 lakh+ users"],
     GOLD),
]
# Timeline line
rect(sl, Inches(0.6), Inches(4.25), Inches(12.1), Inches(0.08), fill=rgb(0x40, 0x55, 0x78))
for i, (phase, period, title, bullets, col) in enumerate(phases):
    bx = Inches(0.3 + i * 3.26)
    by_card = Inches(1.35)
    rect(sl, bx, by_card, Inches(3.05), Inches(4.8), fill=rgb(0x1F, 0x3A, 0x6A))
    rect(sl, bx, by_card, Inches(3.05), Inches(0.07), fill=col)
    # phase badge
    pill(sl, bx + Inches(0.15), by_card + Inches(0.12), Inches(1.0), Inches(0.28),
         fill=col, text=phase, tsize=9)
    txt(sl, period, bx + Inches(1.25), by_card + Inches(0.1), Inches(1.7), Inches(0.36),
        size=10, color=rgb(0xB0, 0xC4, 0xDE), align=PP_ALIGN.RIGHT)
    txt(sl, title, bx + Inches(0.15), by_card + Inches(0.48), Inches(2.75), Inches(0.55),
        size=13, bold=True, color=col)
    for bi, b in enumerate(bullets):
        txt(sl, "▸  " + b, bx + Inches(0.18), by_card + Inches(1.1 + bi * 0.88),
            Inches(2.7), Inches(0.82), size=11, color=rgb(0xCC, 0xD6, 0xE8))
    # timeline dot
    circ = circle(sl, bx + Inches(1.52), Inches(4.29), Inches(0.15), col)

# bottom milestone strip
rect(sl, 0, Inches(6.7), W, Inches(0.8), fill=rgb(0x1F, 0x3A, 0x6A))
for i, (label, val) in enumerate([
    ("Month 6",  "500 insured trips/month"),
    ("Month 12", "₹1 Cr revenue run-rate"),
    ("Month 18", "Pan-India live"),
    ("Month 30", "Profitability milestone"),
    ("Year 3",   "₹15–25 Cr revenue"),
]):
    bx = Inches(0.2 + i * 2.62)
    txt(sl, label, bx, Inches(6.72), Inches(2.5), Inches(0.3),
        size=9, bold=True, color=ORANGE)
    txt(sl, val,   bx, Inches(7.02), Inches(2.5), Inches(0.38),
        size=9.5, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — WHY INVEST
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=LGRAY)
rect(sl, 0, 0, W, Inches(1.3), fill=NAVY)
rect(sl, 0, 0, W, Inches(0.07), fill=ORANGE)

section_label(sl, "11  WHY INVEST")
txt(sl, "5 Compelling Reasons to Bet on Shiftzy Go",
    Inches(0.4), Inches(0.28), Inches(12.5), Inches(0.85),
    size=26, bold=True, color=WHITE)

reasons = [
    ("01", "ZERO COMPETITION", BLUE,
     "No existing platform in India solves peer-to-peer vehicle shifting + travel cost-sharing together. First-mover advantage in a ₹85,000 Cr market."),
    ("02", "ASSET-LIGHT MODEL", TEAL,
     "We own no vehicles. No fleet. No drivers. Pure marketplace — highly scalable with minimal capital, similar to Airbnb and Uber in their early days."),
    ("03", "DUAL-SIDED NETWORK EFFECT", ORANGE,
     "Every new vehicle owner attracts travelers; every new traveler attracts owners. Network effects compound growth automatically."),
    ("04", "MULTIPLE REVENUE STREAMS", GREEN,
     "Platform fees + insurance commissions + premium listings + corporate accounts + data analytics. Diversified from Day 1."),
    ("05", "STRONG SOCIAL IMPACT", GOLD,
     "Reduces individual carbon footprint, makes vehicle ownership costs bearable, creates income opportunity for travelers. ESG-aligned investment."),
]
for i, (num, title, col, desc) in enumerate(reasons):
    by = Inches(1.5 + i * 1.08)
    rect(sl, Inches(0.35), by, Inches(12.63), Inches(0.98), fill=WHITE)
    rect(sl, Inches(0.35), by, Inches(0.07), Inches(0.98), fill=col)
    # number circle
    sh = circle(sl, Inches(0.82), by + Inches(0.49), Inches(0.35), col)
    txt(sl, num, Inches(0.5), by + Inches(0.12), Inches(0.65), Inches(0.72),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, title, Inches(1.3), by + Inches(0.12), Inches(2.5), Inches(0.38),
        size=13, bold=True, color=col)
    txt(sl, desc, Inches(3.9), by + Inches(0.12), Inches(8.9), Inches(0.75),
        size=11, color=DGRAY)
    # arrow
    txt(sl, "→", Inches(1.3), by + Inches(0.5), Inches(0.4), Inches(0.35),
        size=14, bold=True, color=col)

# bottom call out
rect(sl, Inches(0.35), Inches(6.75), Inches(12.63), Inches(0.62), fill=NAVY)
txt(sl, "💡  We are seeking ₹2 – 5 Crore in Seed / Pre-Series A funding to fuel Phase 1 & 2 growth",
    Inches(0.45), Inches(6.8), Inches(12.4), Inches(0.52),
    size=15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — NEXT STEPS / CTA
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=NAVY)
rect(sl, W - Inches(0.55), 0, Inches(0.55), H, fill=ORANGE)
rect(sl, 0, 0, W, Inches(0.07), fill=ORANGE)

section_label(sl, "12  NEXT STEPS")
txt(sl, "Let's Build This Together",
    Inches(0.5), Inches(0.35), Inches(10), Inches(0.9),
    size=34, bold=True, color=WHITE)
txt(sl, "Shiftzy Go is ready. The market is ready. Are you?",
    Inches(0.5), Inches(1.18), Inches(10), Inches(0.55),
    size=18, italic=True, color=ORANGE)

steps_data = [
    ("01", "Sign NDA",            "Confidentiality agreement to share detailed business & financial data"),
    ("02", "Due Diligence",       "Open books — financials, tech, user data, legal structure"),
    ("03", "Term Sheet",          "Agree principal commercial terms (for investors or insurance partner)"),
    ("04", "Pilot Programme",     "100-trip pilot on Chennai–Bangalore to validate model & claims"),
    ("05", "Agreement Signing",   "Formal investment, distribution, or insurance partnership agreement"),
    ("06", "Commercial Launch",   "Go live across 7 South India cities with full insurance coverage"),
]
for i, (num, title, desc) in enumerate(steps_data):
    col_i = i % 2
    row_i = i // 2
    bx = Inches(0.4 + col_i * 6.3)
    by = Inches(1.95 + row_i * 1.58)
    rect(sl, bx, by, Inches(5.9), Inches(1.35), fill=rgb(0x1F, 0x3A, 0x6A))
    rect(sl, bx, by, Inches(5.9), Inches(0.07), fill=ORANGE if col_i == 0 else TEAL)
    sh = circle(sl, bx + Inches(0.42), by + Inches(0.68), Inches(0.32),
                ORANGE if col_i == 0 else TEAL)
    txt(sl, num, bx + Inches(0.12), by + Inches(0.35), Inches(0.62), Inches(0.62),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, title, bx + Inches(0.9), by + Inches(0.14), Inches(4.8), Inches(0.45),
        size=14, bold=True, color=WHITE)
    txt(sl, desc,  bx + Inches(0.9), by + Inches(0.6),  Inches(4.8), Inches(0.65),
        size=11, color=rgb(0xB0, 0xC4, 0xDE))

# Contact box
rect(sl, Inches(0.4), Inches(6.62), Inches(12.13), Inches(0.73), fill=ORANGE)
txt(sl, "📧  business@shiftzygo.com   |   📱  +91 XXXXX XXXXX   |   🌐  www.shiftzygo.com",
    Inches(0.5), Inches(6.7), Inches(12.0), Inches(0.57),
    size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, 0, 0, W, H, fill=NAVY)
rect(sl, 0, 0, W, Inches(0.07), fill=ORANGE)
rect(sl, W - Inches(0.07), 0, Inches(0.07), H, fill=ORANGE)
rect(sl, 0, H - Inches(0.07), W, Inches(0.07), fill=ORANGE)
rect(sl, 0, 0, Inches(0.07), H, fill=ORANGE)

txt(sl, "Thank You",
    Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.8),
    size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(sl, Inches(3.5), Inches(3.2), Inches(6.3), Inches(0.06), fill=ORANGE)

txt(sl, "Shiftzy Go  —  Shift Smart. Travel Free.",
    Inches(1.5), Inches(3.4), Inches(10.3), Inches(0.75),
    size=22, italic=True, color=ORANGE, align=PP_ALIGN.CENTER)

txt(sl, "India's Smart Vehicle Shifting & Travel Cost-Sharing Platform",
    Inches(1.5), Inches(4.15), Inches(10.3), Inches(0.55),
    size=16, color=rgb(0xB0, 0xC4, 0xDE), align=PP_ALIGN.CENTER)

txt(sl, "© 2026 Shiftzy Go Technologies Pvt. Ltd.  ·  Confidential",
    Inches(1.5), Inches(6.8), Inches(10.3), Inches(0.45),
    size=11, color=MGRAY, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "Shiftzy_Go_Investor_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
